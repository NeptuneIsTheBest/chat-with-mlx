import argparse
import atexit
import copy
import hashlib
from pathlib import Path
from typing import Dict, List, Tuple, Optional, Union

import gradio as gr
import pandas as pd
from gradio.components.chatbot import ChatMessage

from .chat import ChatSystemPromptBlock, LoadModelBlock, AdvancedSettingBlock, RAGSettingBlock
from .language import get_text
from .model import Message, MessageRole, ModelManager, Model, VisionModel
from .model_management import AddModelBlock

model_manager = ModelManager()

chat_system_prompt_block = ChatSystemPromptBlock(model_manager=model_manager)
chat_load_model_block = LoadModelBlock(model_manager=model_manager)
chat_advanced_setting_block = AdvancedSettingBlock()
chat_rag_setting_block = RAGSettingBlock()

completion_load_model_block = LoadModelBlock(model_manager=model_manager)

completion_advanced_setting_block = AdvancedSettingBlock()

model_management_add_model_block = AddModelBlock(model_manager=model_manager)


def get_loaded_model() -> Union[Model, VisionModel]:
    model = model_manager.get_loaded_model()
    if model is None:
        raise RuntimeError("No model loaded.")
    return model


def get_file_md5(file_name: Path) -> str:
    md5 = hashlib.md5()
    with open(file_name, "rb") as f:
        for chunk in iter(lambda: f.read(4096), b""):
            md5.update(chunk)
    return md5.hexdigest()


class FileManager:
    def __init__(self):
        self.files = {}

    def format_content(self, file_name, content):
        boundary_start = f"<<<BEGIN FILE:{file_name}>>>"
        boundary_end = f"<<<END FILE:{file_name}>>>"
        content = content.replace(boundary_start, '')
        content = content.replace(boundary_end, '')
        return f"{boundary_start}\n{content}\n{boundary_end}"

    def load_file(self, file_name: Path):
        if file_name.name in self.files:
            if self.files[file_name.name]["md5"] == get_file_md5(file_name):
                return self.files[file_name.name]["content"]
        suffix = file_name.suffix.lower()
        if suffix == ".pdf":
            return self.load_pdf(file_name)
        elif suffix in [".txt", ".csv", ".md"]:
            return self.load_txt_like(file_name)
        elif suffix == ".docx":
            return self.load_docx(file_name)
        elif suffix == ".pptx":
            return self.load_pptx(file_name)
        elif suffix in [".xlsx", ".xls"]:
            return self.load_excel(file_name)
        else:
            return None

    def load_pdf(self, file_name: Path):
        import pypdf
        pdf = pypdf.PdfReader(file_name)
        content = ''
        for page in pdf.pages:
            text = page.extract_text()
            if text:
                content += text
        formatted_content = self.format_content(file_name, content)
        self.files[file_name.name] = {
            "md5": get_file_md5(file_name),
            "content": formatted_content
        }
        return formatted_content

    def load_txt_like(self, file_name: Path):
        with open(file_name, 'r', encoding='utf-8') as f:
            content = f.read()
        formatted_content = self.format_content(file_name, content)
        self.files[file_name.name] = {
            "md5": get_file_md5(file_name),
            "content": formatted_content
        }
        return formatted_content

    def load_docx(self, file_name: Path):
        import docx
        doc = docx.Document(str(file_name))
        content = "\n".join([para.text for para in doc.paragraphs])
        formatted_content = self.format_content(file_name, content)
        self.files[file_name.name] = {
            "md5": get_file_md5(file_name),
            "content": formatted_content
        }
        return formatted_content

    def load_pptx(self, file_name: Path):
        from pptx import Presentation
        prs = Presentation(str(file_name))
        content = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"
        formatted_content = self.format_content(file_name, content)
        self.files[file_name.name] = {
            "md5": get_file_md5(file_name),
            "content": formatted_content
        }
        return formatted_content

    def load_excel(self, file_name: Path):
        import pandas as pd
        excel_file = pd.ExcelFile(file_name)
        content = ""
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            content += f"Sheet: {sheet_name}\n"
            content += df.to_csv(index=False)
        formatted_content = self.format_content(file_name, content)
        self.files[file_name.name] = {
            "md5": get_file_md5(file_name),
            "content": formatted_content
        }
        return formatted_content


file_manager = FileManager()


def preprocess_file(message: Dict, history: List[Dict]) -> Tuple[str, List[Dict]]:
    processed_message = ""
    if "files" in message:
        for file in message["files"]:
            file_content = file_manager.load_file(Path(file))
            processed_message += file_content if file_content else ""
    processed_message += message["text"]
    preprocessed_history = []
    i = 0
    while i < len(history):
        current_history = copy.deepcopy(history[i])
        if isinstance(current_history["content"], tuple):
            for file in current_history["content"]:
                file_content = file_manager.load_file(Path(file))
                current_history["content"] = file_content if file_content else ""
                if i + 1 < len(history):
                    next_history = copy.deepcopy(history[i + 1])
                    current_history["content"] += next_history["content"] if next_history else ""
                    i += 1
                preprocessed_history.append(current_history)
                current_history = copy.deepcopy(history[i])
            i += 1
            continue
        preprocessed_history.append(current_history)
        i += 1
    return processed_message, preprocessed_history


def handle_chat(message: Dict,
                history: List[Dict],
                system_prompt: str = None,
                temperature: float = 0.7,
                top_p: float = 0.9,
                max_tokens: int = 512,
                repetition_penalty: float = 1.0,
                stream: bool = True):
    try:
        model = get_loaded_model()
        if isinstance(model, VisionModel):
            if model.processor.tokenizer.chat_template is None:
                raise RuntimeError("No chat template.")
        else:
            if model.tokenizer.chat_template is None:
                raise RuntimeError("No chat template.")

        if system_prompt and system_prompt.strip() != "":
            history = [Message(MessageRole.SYSTEM, content=system_prompt).to_dict()] + history

        images = []
        if isinstance(model, VisionModel):
            for h in history:
                if "content" in h:
                    for file in h["content"]:
                        if Path(file).suffix.lower() in [".jpg", ".png", ".jpeg"]:
                            images.append(file)
            if "files" in message:
                for file in message["files"]:
                    if Path(file).suffix.lower() in [".jpg", ".png", ".jpeg"]:
                        images.append(file)

        message, history = preprocess_file(message, history)

        temperature = float(temperature)
        top_p = float(top_p)
        repetition_penalty = float(repetition_penalty)

        response = ChatMessage(role="assistant", content="")
        if isinstance(model, VisionModel):
            eos_token = model.processor.tokenizer.eos_token
            for chunk in model.generate_response(
                    message=message,
                    images=images,
                    history=history,
                    stream=stream,
                    temperature=temperature,
                    top_p=top_p,
                    max_tokens=max_tokens,
                    repetition_penalty=repetition_penalty):
                if eos_token not in chunk:
                    response.content += chunk
                    yield response
        else:
            eos_token = model.tokenizer.eos_token
            for chunk in model.generate_response(
                    message=message,
                    history=history,
                    stream=stream,
                    temperature=temperature,
                    top_p=top_p,
                    max_tokens=max_tokens,
                    repetition_penalty=repetition_penalty):
                if eos_token not in chunk.text:
                    response.content += chunk.text
                    yield response
    except Exception as e:
        raise gr.Error(str(e))


def handle_completion(prompt: str,
                      temperature: float = 0.7,
                      top_p: float = 0.9,
                      max_tokens: int = 512,
                      repetition_penalty: float = 1.0,
                      stream: bool = True):
    try:
        model = get_loaded_model()
        if isinstance(model, VisionModel):
            raise RuntimeError("Not supported yet.")

        temperature = float(temperature)
        top_p = float(top_p)
        repetition_penalty = float(repetition_penalty)

        if not stream:
            return prompt + model.generate_completion(
                prompt=prompt,
                stream=stream,
                temperature=temperature,
                top_p=top_p,
                max_tokens=max_tokens,
                repetition_penalty=repetition_penalty)
        else:
            response = prompt
            eos_token = model.tokenizer.eos_token
            for chunk in model.generate_completion(
                    prompt=prompt,
                    stream=stream,
                    temperature=temperature,
                    top_p=top_p,
                    max_tokens=max_tokens,
                    repetition_penalty=repetition_penalty):
                if eos_token not in chunk.text:
                    response += chunk.text
                    yield response
    except Exception as e:
        raise gr.Error(str(e))


def get_load_model_status():
    if model_manager.get_loaded_model_config():
        return get_text("Page.Chat.LoadModelBlock.Textbox.model_status.loaded_value").format(model_manager.get_loaded_model_config().get("display_name"))
    else:
        return get_text("Page.Chat.LoadModelBlock.Textbox.model_status.not_loaded_value")


def load_model(model_name: str) -> Tuple[str, str]:
    try:
        model_manager.load_model(model_name)
        return get_load_model_status(), model_manager.get_system_prompt(default=True)
    except Exception as e:
        raise gr.Error(str(e))


def chat_load_model_callback(model_name: str):
    return load_model(model_name)


def completion_load_model_callback(model_name: str):
    return load_model(model_name)[0]


def get_default_system_prompt_callback():
    if model_manager.get_loaded_model_config():
        return model_manager.get_system_prompt(default=True)
    else:
        raise gr.Error("No model loaded.")


def update_model_management_models_list():
    return pd.DataFrame({get_text("Page.ModelManagement.Dataframe.model_list.headers"): model_manager.get_model_list()})


def update_model_selector_choices():
    return gr.update(choices=model_manager.get_model_list(), value=chat_load_model_block.update_select_model_dropdown_value())


def add_model(model_name: Optional[str], original_repo: str, mlx_repo: str, quantize: str, default_language: str, default_system_prompt: Optional[str], multimodal_ability: List[str]):
    try:
        model_manager.add_config(original_repo, mlx_repo, model_name, quantize, default_language, default_system_prompt, multimodal_ability)
    except Exception as e:
        raise gr.Error(str(e))


with gr.Blocks(fill_height=True, fill_width=True, title="Chat with MLX") as app:
    gr.HTML("<h1>Chat with MLX</h1>")

    with gr.Tab(get_text("Tab.chat")):
        with gr.Row():
            with gr.Column(scale=2):
                with gr.Row():
                    gr.Markdown(f"## {get_text('Page.Chat.Markdown.configuration')}")

                    chat_load_model_block.render_all()
                    chat_load_model_block.model_selector_dropdown.select(
                        fn=lambda x: x,
                        inputs=[chat_load_model_block.model_selector_dropdown],
                        outputs=[completion_load_model_block.model_selector_dropdown]
                    )
                    chat_load_model_block.load_model_button.click(
                        fn=chat_load_model_callback,
                        inputs=[chat_load_model_block.model_selector_dropdown],
                        outputs=[
                            chat_load_model_block.model_status_textbox,
                            chat_system_prompt_block.system_prompt_textbox
                        ]
                    ).then(
                        fn=lambda x: x,
                        inputs=[chat_load_model_block.model_status_textbox],
                        outputs=[completion_load_model_block.model_status_textbox]
                    )

                with gr.Accordion(label=get_text("Page.Chat.Accordion.AdvancedSetting.label"), open=False):
                    chat_advanced_setting_block.render_all()

                with gr.Accordion(label=get_text("Page.Chat.Accordion.RAGSetting.label"), open=False):
                    chat_rag_setting_block.render_all()

            with gr.Column(scale=8):
                with gr.Row(equal_height=True):
                    chat_system_prompt_block.render_all()
                    chat_system_prompt_block.default_system_prompt_button.click(
                        fn=get_default_system_prompt_callback,
                        outputs=[chat_system_prompt_block.system_prompt_textbox]
                    )
                    chat_system_prompt_block.system_prompt_textbox.change(
                        fn=model_manager.set_custom_prompt,
                        inputs=[chat_system_prompt_block.system_prompt_textbox]
                    )

                chatbot = gr.Chatbot(
                    type="messages",
                    show_copy_button=True,
                    render=False,
                    latex_delimiters=[
                        {"left": "$$", "right": "$$", "display": True},
                        {"left": "$", "right": "$", "display": False},
                        {"left": "\\(", "right": "\\)", "display": False},
                        {"left": "\\begin{equation}", "right": "\\end{equation}", "display": True},
                        {"left": "\\begin{align}", "right": "\\end{align}", "display": True},
                        {"left": "\\begin{alignat}", "right": "\\end{alignat}", "display": True},
                        {"left": "\\begin{gather}", "right": "\\end{gather}", "display": True},
                        {"left": "\\begin{CD}", "right": "\\end{CD}", "display": True},
                        {"left": "\\[", "right": "\\]", "display": True}
                    ]
                )

                gr.ChatInterface(
                    multimodal=True,
                    chatbot=chatbot,
                    type="messages",
                    fn=handle_chat,
                    title=None,
                    autofocus=False,
                    additional_inputs=[
                        chat_system_prompt_block.system_prompt_textbox,
                        chat_advanced_setting_block.temperature_slider,
                        chat_advanced_setting_block.top_p_slider,
                        chat_advanced_setting_block.max_tokens_slider,
                        chat_advanced_setting_block.repetition_penalty_slider
                    ]
                )

    with gr.Tab(get_text("Tab.completion"), interactive=True):
        with gr.Row():
            with gr.Column(scale=2):
                gr.Markdown(f"## {get_text('Page.Chat.Markdown.configuration')}")

                completion_load_model_block.render_all()
                completion_load_model_block.model_selector_dropdown.select(
                    fn=lambda x: x,
                    inputs=[completion_load_model_block.model_selector_dropdown],
                    outputs=[chat_load_model_block.model_selector_dropdown]
                )
                completion_load_model_block.load_model_button.click(
                    fn=completion_load_model_callback,
                    inputs=[completion_load_model_block.model_selector_dropdown],
                    outputs=[
                        completion_load_model_block.model_status_textbox
                    ]
                ).then(
                    fn=lambda x: x,
                    inputs=[completion_load_model_block.model_status_textbox],
                    outputs=[chat_load_model_block.model_status_textbox]
                )

                with gr.Row(visible=False):
                    with gr.Accordion(label=get_text("Page.Chat.Accordion.AdvancedSetting.label"), open=True):
                        completion_advanced_setting_block.render_all()

            with gr.Column(scale=8):
                completion_textbox = gr.Textbox(lines=25, render=False)
                completion_interface = gr.Interface(
                    clear_btn=None,
                    flagging_mode="never",
                    fn=handle_completion,
                    inputs=[
                        gr.Textbox(lines=10, show_copy_button=True, render=True, label=get_text("Page.Completion.Textbox.prompt.label")),
                        completion_advanced_setting_block.temperature_slider,
                        completion_advanced_setting_block.top_p_slider,
                        completion_advanced_setting_block.max_tokens_slider,
                        completion_advanced_setting_block.repetition_penalty_slider
                    ],
                    outputs=[
                        gr.Textbox(lines=25, show_copy_button=True, render=True, label=get_text("Page.Completion.Textbox.output.label"))
                    ],
                    submit_btn=get_text("Page.Completion.Button.submit.value"),
                    stop_btn=get_text("Page.Completion.Button.stop.value"),
                )

    model_list = gr.Dataframe(
        headers=[get_text("Page.ModelManagement.Dataframe.model_list.headers")],
        value=update_model_management_models_list(),
        datatype=["str"],
        row_count=(10, "dynamic"),
        render=False,
        interactive=False
    )

    with gr.Tab(get_text("Tab.model_management"), interactive=True):
        with gr.Row(equal_height=True):
            with gr.Column(scale=5):
                model_list.render()

            with gr.Column(scale=5):
                model_management_add_model_block.render_all()
                model_management_add_model_block.add_button.click(
                    fn=add_model,
                    inputs=[
                        model_management_add_model_block.model_name_textbox,
                        model_management_add_model_block.original_repo_textbox,
                        model_management_add_model_block.mlx_repo_textbox,
                        model_management_add_model_block.quantize_dropdown,
                        model_management_add_model_block.default_language_dropdown,
                        model_management_add_model_block.default_system_prompt_textbox,
                        model_management_add_model_block.multimodal_ability_dropdown
                    ]
                ).then(
                    fn=update_model_management_models_list,
                    outputs=[
                        model_list
                    ]
                ).then(
                    fn=update_model_selector_choices,
                    outputs=[
                        chat_load_model_block.model_selector_dropdown
                    ]
                ).then(
                    fn=update_model_selector_choices,
                    outputs=[
                        completion_load_model_block.model_selector_dropdown
                    ]
                )

    app.load(
        fn=update_model_management_models_list,
        outputs=[
            model_list
        ]
    ).then(
        fn=update_model_selector_choices,
        outputs=[
            chat_load_model_block.model_selector_dropdown
        ]
    ).then(
        fn=update_model_selector_choices,
        outputs=[
            completion_load_model_block.model_selector_dropdown
        ]
    )


def exit_handler():
    if model_manager.get_loaded_model():
        model_manager.close_model()


atexit.register(exit_handler)


def start(port: int, share: bool = False, in_browser: bool = True) -> None:
    print(f"Starting the app on port {port} with share={share} and in_browser={in_browser}")
    app.launch(server_port=port, inbrowser=in_browser, share=share)


def main():
    parser = argparse.ArgumentParser(description="Chat with MLX")
    parser.add_argument(
        "--port",
        type=int,
        default=7860,
        help="The port number to run the application on (default: 7860)"
    )
    parser.add_argument(
        "--share",
        action="store_true",
        help="Enable sharing the application link externally"
    )
    parser.add_argument(
        "--no-browser",
        action="store_true",
        help="Do not open the application in the default web browser"
    )
    args = parser.parse_args()

    start(port=args.port, share=args.share, in_browser=not args.no_browser)


if __name__ == "__main__":
    main()
